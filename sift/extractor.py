"""
Content extraction module for Sift.
Extracts text content from various document formats.
Supports PDF vision fallback and legacy .doc files.
"""

import base64
import os
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

from .config import SiftConfig


class ContentExtractor:
    """Extracts text content from documents."""

    # Minimum chars to consider text extraction successful
    MIN_EXTRACTED_CHARS = 50

    def __init__(self, config: SiftConfig):
        self.config = config
        self._docx_available = False
        self._pdf_available = False
        self._xlsx_available = False
        self._pptx_available = False
        self._pdf2image_available = False
        self._pillow_available = False
        self._textract_available = False
        self._antiword_available = False

        # Try importing optional dependencies
        try:
            import docx
            self._docx_available = True
        except ImportError:
            pass

        try:
            import PyPDF2
            self._pdf_available = True
        except ImportError:
            pass

        try:
            import openpyxl
            self._xlsx_available = True
        except ImportError:
            pass

        try:
            import pptx
            self._pptx_available = True
        except ImportError:
            pass

        try:
            from pdf2image import convert_from_path
            self._pdf2image_available = True
        except ImportError:
            pass

        try:
            from PIL import Image
            self._pillow_available = True
        except ImportError:
            pass

        # Check for antiword (Windows/Linux tool for .doc files)
        try:
            result = subprocess.run(
                ["antiword", "--version"],
                capture_output=True,
                timeout=5
            )
            self._antiword_available = True
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # Initialize client lazily for vision fallback
        self._client = None

    def _get_client(self):
        """Get Anthropic client lazily."""
        if self._client is None:
            self._client = self.config.get_anthropic_client()
        return self._client

    def extract(self, file_path: Path, max_chars: int = 5000) -> Optional[str]:
        """
        Extract text content from a file.
        Returns None if extraction fails or file type not supported.
        """
        if not file_path.exists():
            return None

        # Check file size
        if file_path.stat().st_size > self.config.max_file_size_bytes():
            return f"[File too large for content extraction: {file_path.stat().st_size / 1024 / 1024:.1f} MB]"

        ext = file_path.suffix.lower()

        try:
            if ext == ".txt":
                return self._extract_txt(file_path, max_chars)
            elif ext == ".csv":
                return self._extract_csv(file_path, max_chars)
            elif ext == ".docx":
                return self._extract_docx(file_path, max_chars)
            elif ext == ".doc":
                return self._extract_doc_legacy(file_path, max_chars)
            elif ext == ".pdf":
                return self._extract_pdf(file_path, max_chars)
            elif ext in (".xlsx", ".xls"):
                return self._extract_xlsx(file_path, max_chars)
            elif ext in (".pptx", ".ppt"):
                return self._extract_pptx(file_path, max_chars)
            elif ext == ".rtf":
                return self._extract_rtf(file_path, max_chars)
            else:
                return None
        except Exception as e:
            return f"[Content extraction failed: {str(e)}]"

    def _extract_txt(self, file_path: Path, max_chars: int) -> str:
        """Extract text from plain text file."""
        encodings = ["utf-8", "utf-16", "latin-1", "cp1252"]

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    content = f.read(max_chars)
                    return content
            except (UnicodeDecodeError, UnicodeError):
                continue

        return "[Unable to decode text file]"

    def _extract_csv(self, file_path: Path, max_chars: int) -> str:
        """Extract content from CSV file (first few rows as sample)."""
        content = self._extract_txt(file_path, max_chars)
        if content.startswith("["):
            return content

        # Return first portion as sample
        lines = content.split("\n")[:20]
        return "\n".join(lines)

    def _extract_docx(self, file_path: Path, max_chars: int) -> Optional[str]:
        """Extract text from Word document (.docx)."""
        if not self._docx_available:
            return "[python-docx not installed - cannot extract .docx content]"

        import docx

        try:
            doc = docx.Document(str(file_path))
            paragraphs = []
            char_count = 0

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
                    char_count += len(text)
                    if char_count >= max_chars:
                        break

            return "\n\n".join(paragraphs)[:max_chars]
        except Exception as e:
            return f"[DOCX extraction failed: {str(e)}]"

    def _extract_doc_legacy(self, file_path: Path, max_chars: int) -> Optional[str]:
        """Extract text from legacy Word document (.doc) using antiword or COM."""
        # Method 1: Try antiword (cross-platform)
        if self._antiword_available:
            try:
                result = subprocess.run(
                    ["antiword", str(file_path)],
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                if result.returncode == 0 and result.stdout.strip():
                    return result.stdout[:max_chars]
            except (subprocess.TimeoutExpired, Exception):
                pass

        # Method 2: Try Windows COM automation (Word must be installed)
        if os.name == 'nt':
            try:
                import win32com.client
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                try:
                    doc = word.Documents.Open(str(file_path), ReadOnly=True)
                    text = doc.Content.Text
                    doc.Close(False)
                    return text[:max_chars]
                finally:
                    word.Quit()
            except ImportError:
                pass  # pywin32 not installed
            except Exception as e:
                pass  # Word not available or file error

        # Method 3: Try textract as last resort
        try:
            import textract
            text = textract.process(str(file_path)).decode('utf-8', errors='ignore')
            return text[:max_chars]
        except ImportError:
            pass
        except Exception:
            pass

        return "[Legacy .doc extraction failed - install antiword, pywin32, or textract]"

    def _extract_pdf(self, file_path: Path, max_chars: int) -> Optional[str]:
        """Extract text from PDF file with vision fallback."""
        text = None

        # Try text extraction first
        if self._pdf_available:
            import PyPDF2

            try:
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    text_parts = []
                    char_count = 0

                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
                            char_count += len(page_text)
                            if char_count >= max_chars:
                                break

                    text = "\n\n".join(text_parts)[:max_chars]
            except Exception:
                pass

        # Check if extraction was successful (got meaningful content)
        if text and len(text.strip()) >= self.MIN_EXTRACTED_CHARS:
            return text

        # Fallback to vision-based extraction
        vision_text = self._extract_pdf_with_vision(file_path, max_chars)
        if vision_text:
            return vision_text

        # Return whatever we got, or error message
        if text:
            return text
        return "[PDF extraction failed - text extraction returned empty and vision fallback unavailable]"

    def _extract_pdf_with_vision(self, file_path: Path, max_chars: int) -> Optional[str]:
        """Extract text from PDF using Claude Vision API."""
        if not self._pdf2image_available or not self._pillow_available:
            return None

        try:
            from pdf2image import convert_from_path
            from PIL import Image
            import io

            # Convert first 3 pages to images
            images = convert_from_path(
                str(file_path),
                first_page=1,
                last_page=3,
                dpi=150  # Balance quality vs size
            )

            if not images:
                return None

            # Convert images to base64
            image_contents = []
            for img in images:
                # Resize if too large (max 1568px on longest side per Anthropic docs)
                max_dim = 1568
                if max(img.size) > max_dim:
                    ratio = max_dim / max(img.size)
                    new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
                    img = img.resize(new_size, Image.Resampling.LANCZOS)

                # Convert to PNG bytes
                buffer = io.BytesIO()
                img.save(buffer, format="PNG", optimize=True)
                img_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')

                image_contents.append({
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": "image/png",
                        "data": img_base64
                    }
                })

            # Add text prompt
            image_contents.append({
                "type": "text",
                "text": "Extract all text content from this document. Return the full text, preserving the structure and formatting as much as possible. If this is an invoice or form, include all field labels and values."
            })

            # Call Claude Vision
            client = self._get_client()
            response = client.messages.create(
                model=self.config.model,
                max_tokens=4096,
                messages=[{
                    "role": "user",
                    "content": image_contents
                }]
            )

            return response.content[0].text[:max_chars]

        except Exception as e:
            return None

    def _extract_xlsx(self, file_path: Path, max_chars: int) -> Optional[str]:
        """Extract content from Excel file."""
        if not self._xlsx_available:
            return "[openpyxl not installed - cannot extract .xlsx content]"

        import openpyxl

        try:
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            content_parts = []
            char_count = 0

            for sheet_name in wb.sheetnames[:5]:  # First 5 sheets
                sheet = wb[sheet_name]
                content_parts.append(f"--- Sheet: {sheet_name} ---")

                row_count = 0
                for row in sheet.iter_rows(max_row=50, values_only=True):  # First 50 rows
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        content_parts.append(row_text)
                        char_count += len(row_text)
                        row_count += 1

                    if char_count >= max_chars:
                        break

                if char_count >= max_chars:
                    break

            wb.close()
            return "\n".join(content_parts)[:max_chars]
        except Exception as e:
            return f"[XLSX extraction failed: {str(e)}]"

    def _extract_pptx(self, file_path: Path, max_chars: int) -> Optional[str]:
        """Extract text from PowerPoint file."""
        if not self._pptx_available:
            return "[python-pptx not installed - cannot extract .pptx content]"

        from pptx import Presentation

        try:
            prs = Presentation(str(file_path))
            content_parts = []
            char_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    content_parts.append(f"--- Slide {slide_num} ---")
                    content_parts.extend(slide_text)
                    char_count += sum(len(t) for t in slide_text)

                if char_count >= max_chars:
                    break

            return "\n".join(content_parts)[:max_chars]
        except Exception as e:
            return f"[PPTX extraction failed: {str(e)}]"

    def _extract_rtf(self, file_path: Path, max_chars: int) -> Optional[str]:
        """Extract text from RTF file (basic extraction)."""
        try:
            with open(file_path, "rb") as f:
                content = f.read(max_chars * 2)  # RTF has overhead

            # Very basic RTF stripping - just remove RTF commands
            text = content.decode("latin-1", errors="ignore")

            # Remove RTF header/commands (basic approach)
            import re
            text = re.sub(r"\\[a-z]+\d*\s?", " ", text)
            text = re.sub(r"[{}]", "", text)
            text = re.sub(r"\s+", " ", text).strip()

            return text[:max_chars]
        except Exception as e:
            return f"[RTF extraction failed: {str(e)}]"

    def get_extraction_capabilities(self) -> dict[str, bool]:
        """Return which file types can be extracted."""
        # .doc is available if antiword, pywin32, or textract is available
        doc_available = self._antiword_available or os.name == 'nt'  # Windows has COM fallback

        return {
            ".txt": True,
            ".csv": True,
            ".docx": self._docx_available,
            ".doc": doc_available,
            ".pdf": self._pdf_available or (self._pdf2image_available and self._pillow_available),
            ".pdf_vision": self._pdf2image_available and self._pillow_available,
            ".xlsx": self._xlsx_available,
            ".xls": self._xlsx_available,
            ".pptx": self._pptx_available,
            ".ppt": self._pptx_available,
            ".rtf": True,
        }
